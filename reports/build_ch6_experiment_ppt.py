from __future__ import annotations

import csv
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import pptx
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = ROOT / "thesis" / "figures"
CSV_DIR = ROOT / "code" / "python" / "experiments" / "results"
REPORT_DIR = ROOT / "reports" / "ch6_experiment_report"
HTML_PATH = REPORT_DIR / "index.html"
PPTX_PATH = REPORT_DIR / "ch6_experiment_talk.pptx"
BUILD_REPORT = REPORT_DIR / "CH6_PPT_BUILD.md"

FONT = "Microsoft YaHei"
TITLE = RGBColor(20, 38, 62)
TEXT = RGBColor(38, 48, 63)
MUTED = RGBColor(91, 106, 125)
BLUE = RGBColor(24, 86, 145)
TEAL = RGBColor(0, 126, 118)
GREEN = RGBColor(35, 145, 91)
GRAY_BG = RGBColor(244, 247, 251)
LINE = RGBColor(212, 222, 234)
WHITE = RGBColor(255, 255, 255)
WARN = RGBColor(158, 91, 0)

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)


@dataclass(frozen=True)
class FigureUse:
    slide: int
    filename: str
    purpose: str


@dataclass(frozen=True)
class CsvUse:
    filename: str
    purpose: str


FIGURES = [
    FigureUse(6, "exp01_convergence.png", "Dirichlet convergence"),
    FigureUse(7, "exp02_temperature_field_comparison.png", "nonhomogeneous Dirichlet visual check"),
    FigureUse(9, "exp03_neumann_mixed_summary.png", "Neumann/mixed boundary summary"),
    FigureUse(11, "exp06_accuracy_cost_error_time.png", "accuracy-cost comparison"),
    FigureUse(12, "exp06_time_scaling.png", "time scaling"),
    FigureUse(14, "exp07_spectral_denominator_heatmaps.png", "small-denominator risk map"),
    FigureUse(15, "exp04_condition_check.png", "condition-number equivalence check"),
    FigureUse(16, "exp04_gmres_history.png", "exp04 GMRES residual history"),
    FigureUse(18, "exp05_near_resonance_summary.png", "baseline near-resonance scan"),
    FigureUse(19, "exp05_multimode_resonance_summary.png", "multimode resonance summary"),
    FigureUse(20, "exp05_dominant_mode_projection.png", "dominant modal projection"),
    FigureUse(21, "exp05_resonance_gmres_history.png", "near-resonance GMRES history"),
]

CSVS = [
    CsvUse("exp00_fft_vs_sparse.csv", "implementation validation"),
    CsvUse("exp01_convergence.csv", "Dirichlet convergence"),
    CsvUse("exp02_nonhom_bc.csv", "nonhomogeneous Dirichlet check"),
    CsvUse("exp03_neumann_mixed.csv", "Neumann/mixed diagnostics"),
    CsvUse("exp04_modified_vs_true.csv", "Helmholtz spectral indicators"),
    CsvUse("exp04_condition_check.csv", "condition check"),
    CsvUse("exp04_gmres_history.csv", "GMRES residual history"),
    CsvUse("exp05_resonance.csv", "near-resonance baseline"),
    CsvUse("exp05_multimode_resonance.csv", "multimode resonance"),
    CsvUse("exp05_resonance_gmres_history.csv", "near-resonance residual history"),
    CsvUse("exp06_accuracy_cost.csv", "accuracy-cost comparison"),
    CsvUse("exp07_spectral_denominator_summary.csv", "spectral denominator maps"),
]


def rel(path: Path) -> str:
    try:
        return str(path.relative_to(ROOT)).replace("\\", "/")
    except ValueError:
        return str(path)


def set_font(run, size: float, bold: bool = False, color: RGBColor = TEXT) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 19,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold=bold, color=color)


def add_multiline(
    slide,
    lines: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: RGBColor = TEXT,
    bullet: bool = True,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = f"• {line}" if bullet else line
        set_font(run, size, color=color)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.62, 0.34, 11.9, 0.52, size=30, bold=True, color=TITLE)
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.92, 11.8, 0.35, size=14, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.20), Inches(12.1), Inches(0.018))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, number: int) -> None:
    add_text(slide, f"{number:02d}", 12.42, 7.06, 0.58, 0.18, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_takeaway(slide, text: str, y: float = 6.67) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(y - 0.10), Inches(11.85), Inches(0.014))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    add_text(slide, f"要点：{text}", 0.72, y, 11.85, 0.42, size=17, bold=True, color=TEAL)


def add_formula(slide, text: str, x: float, y: float, w: float, h: float, size: float = 23) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_BG
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size, bold=True, color=TITLE)


def add_image_fit(slide, image_path: Path, x: float, y: float, w: float, h: float, border: bool = True) -> bool:
    if not image_path.exists():
        ph = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(252, 238, 238)
        ph.line.color.rgb = RGBColor(170, 60, 60)
        add_text(slide, f"缺失图片：{image_path.name}", x + 0.2, y + h / 2 - 0.12, w - 0.4, 0.3, size=17, bold=True, color=RGBColor(150, 45, 45), align=PP_ALIGN.CENTER)
        return False
    with Image.open(image_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
    else:
        draw_h = h
        draw_w = h * img_ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(draw_x), Inches(draw_y), width=Inches(draw_w), height=Inches(draw_h))
    if border:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))
        rect.fill.background()
        rect.line.color.rgb = LINE
    return True


def add_note(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.25), Inches(7.75), Inches(12.6), Inches(0.85))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "SPEAKER_NOTES: " + text
    set_font(run, 8, color=RGBColor(120, 120, 120))


def new_slide(prs: Presentation, idx: int, title: str | None = None, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    if title:
        add_title(slide, title, subtitle)
    add_footer(slide, idx)
    return slide


def read_csv_preview(name: str, max_rows: int = 4) -> list[dict[str, str]]:
    path = CSV_DIR / name
    if not path.exists():
        return []
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return [row for _, row in zip(range(max_rows), csv.DictReader(f))]


def fmt(value: str | float | int | None) -> str:
    if value is None or value == "":
        return "-"
    try:
        x = float(value)
    except (TypeError, ValueError):
        return str(value)
    if x == 0:
        return "0"
    if abs(x) < 1e-3 or abs(x) >= 1e4:
        return f"{x:.2e}"
    return f"{x:.4g}"


def add_table(slide, data: list[list[str]], x: float, y: float, w: float, h: float, col_widths: list[float] | None = None, font_size: float = 11.5) -> None:
    rows = len(data)
    cols = max(len(row) for row in data)
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for col, width in enumerate(col_widths):
            table.columns[col].width = Inches(width)
    for r, row in enumerate(data):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = row[c] if c < len(row) else ""
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(232, 239, 248) if r == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    set_font(run, font_size, bold=(r == 0), color=TITLE if r == 0 else TEXT)


def add_flow(slide) -> None:
    labels = [
        ("1", "Dirichlet\n收敛"),
        ("2", "Neumann/mixed\n边界"),
        ("3", "精度--成本\n对比"),
        ("4", "Helmholtz 谱\n与 GMRES"),
        ("5", "近共振\n模态放大"),
    ]
    x0, y0, gap = 0.78, 2.05, 2.43
    for i, (num, label) in enumerate(labels):
        x = x0 + i * gap
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y0), Inches(1.12), Inches(1.12))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE if i < 3 else TEAL
        circ.line.fill.background()
        tf = circ.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, 27, bold=True, color=WHITE)
        add_text(slide, label, x - 0.35, y0 + 1.32, 1.85, 0.75, size=16, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x + 1.26), Inches(y0 + 0.43), Inches(0.76), Inches(0.26))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LINE
            arrow.line.fill.background()


def build_deck() -> tuple[list[str], int]:
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    missing: list[str] = []
    note_count = 0

    def note(slide, text: str) -> None:
        nonlocal note_count
        note_count += 1
        add_note(slide, text)

    # 1
    slide = new_slide(prs, 1)
    add_text(slide, "第六章数值实验报告", 0.72, 1.30, 11.9, 0.55, size=36, bold=True, color=TITLE)
    add_text(slide, "FFT 快速求解器与无预处理 GMRES 对比", 0.72, 1.92, 11.9, 0.42, size=25, bold=True, color=TEAL)
    add_text(slide, "收敛阶 · 边界处理 · 精度--成本 · Helmholtz 小分母机制", 0.74, 2.55, 11.6, 0.34, size=18, color=MUTED)
    add_formula(slide, "(-∇² + σ)u = f", 0.95, 4.05, 4.7, 0.75, size=25)
    add_formula(slide, "dₚ,q = λₚ,q + σ", 6.1, 4.05, 4.7, 0.75, size=25)
    add_text(slide, "答辩定位：只讲第六章实验设计、证据链和结论边界", 0.8, 6.25, 11.6, 0.32, size=16, color=MUTED)
    note(slide, "开场先说明这不是全文理论推导，而是第六章实验报告。后面所有图都围绕统一方程和频域分母展开。")

    # 2
    slide = new_slide(prs, 2, "实验章主线", "从离散格式到 near-resonance 机制")
    add_flow(slide)
    add_multiline(slide, [
        "前两组实验确认离散格式和边界处理可靠",
        "第三组实验直接回应“FFT 与迭代法对比”",
        "后两组实验解释 true Helmholtz 小分母如何导致模态放大和 GMRES 停滞",
    ], 1.05, 4.75, 11.1, 1.15, size=18)
    add_takeaway(slide, "第六章形成“理论符号--实现路径--数值证据--机制解释”的闭环。")
    note(slide, "这一页给评委一个总地图。强调 exp00 是实现校验，不进入五个核心实验编号。")

    # 3
    slide = new_slide(prs, 3, "统一方程与频域分母", "区分 Poisson、modified Helmholtz 与 true Helmholtz")
    add_formula(slide, "Poisson: σ = 0", 0.9, 1.65, 3.65, 0.70, size=21)
    add_formula(slide, "modified: σ = +κ²", 4.86, 1.65, 3.65, 0.70, size=21)
    add_formula(slide, "true: σ = -κ²", 8.82, 1.65, 3.65, 0.70, size=21)
    add_multiline(slide, [
        "矩形规则网格上，DST/DCT 使离散 Laplacian 在频域逐模态对角化",
        "每个模态的求解可以理解为 ûₚ,q = f̂ₚ,q / dₚ,q",
        "true Helmholtz 中 dₚ,q = λₚ,q - κ²，可能接近 0",
        "小分母是解放大和无预处理 GMRES30 困难的共同原因",
    ], 1.0, 3.05, 11.1, 1.78, size=19)
    add_takeaway(slide, "Helmholtz 实验的核心变量不是图像颜色，而是 dₚ,q 距离 0 有多近。")
    note(slide, "把 modified 和 true 的区别说清楚：modified 分母全正，true 分母可能正负分裂并接近零。")

    # 4
    slide = new_slide(prs, 4, "exp00：实现校验", "FFT solver vs sparse direct")
    exp00 = read_csv_preview("exp00_fft_vs_sparse.csv", 4)
    err = "10⁻¹⁴"
    if exp00:
        for key in ("error_vs_sparse", "linf_error", "max_error"):
            if key in exp00[0]:
                err = fmt(exp00[0][key])
                break
    add_multiline(slide, [
        f"表 5 比较 FFT 解与 sparse direct 解，误差约为 {err} 量级",
        "它只说明两条求解路径对齐到同一个离散线性系统",
        "它不证明连续 PDE 精度，也不替代后续收敛阶实验",
        "FFT9 与五点 sparse system 属于不同离散模板，不混入该校验表",
    ], 0.95, 1.65, 11.3, 2.05, size=20)
    add_table(slide, [["定位", "回答的问题"], ["implementation validation", "代码路径是否对齐"], ["非核心实验", "不回答连续误差阶数"], ["后续实验", "验证精度、成本和机制"]], 2.1, 4.2, 9.0, 1.35, [3.0, 6.0], font_size=12.5)
    add_takeaway(slide, "exp00 是可信实验链路的地基，不是主要科学贡献。")
    note(slide, "答辩时不要把 exp00 讲成精度结果。它只是排除实现路径错误。")

    # 5
    slide = new_slide(prs, 5, "实验一：Dirichlet 收敛目标", "先确认格式阶数，再看边界子情形")
    add_multiline(slide, [
        "五点格式：FA / CR / FACR-like，理论二阶 O(h²)",
        "FFT9：Dirichlet 九点紧致格式，理论四阶 O(h⁴)",
        "polynomial manufactured solution 避免只依赖单一 Fourier 模态",
        "非齐次 Dirichlet 子情形检查边界值移到右端项的符号和模板",
    ], 0.95, 1.55, 11.25, 2.05, size=20)
    add_formula(slide, "h → h/2:  二阶误差约 /4，四阶误差约 /16", 1.6, 4.2, 10.2, 0.80, size=22)
    add_takeaway(slide, "收敛图看斜率，不只看单个误差数字。")
    note(slide, "这一页先讲读图方法。单一 Fourier 模态只作为 sanity check，主要证据是 polynomial manufactured solution。")

    # 6
    slide = new_slide(prs, 6, "实验一：Dirichlet 收敛图", "五点二阶与 FFT9 四阶")
    img = FIG_DIR / "exp01_convergence.png"
    if not add_image_fit(slide, img, 0.72, 1.22, 11.9, 5.23):
        missing.append(img.name)
    add_takeaway(slide, "FA/CR/FACR-like 跟随二阶斜率，FFT9 跟随四阶斜率。")
    note(slide, "讲图时聚焦曲线斜率和参考线。不要把这张图扩展成所有边界条件下 FFT9 都四阶。")

    # 7
    slide = new_slide(prs, 7, "非齐次 Dirichlet 子情形", "边界修正项的 visual sanity check")
    img = FIG_DIR / "exp02_temperature_field_comparison.png"
    if not add_image_fit(slide, img, 0.72, 1.20, 11.9, 5.22):
        missing.append(img.name)
    add_takeaway(slide, "边界值移至右端项的符号和模板处理正确；该图不是独立核心实验。")
    note(slide, "讲 exact、numerical、error 的对应关系。强调边界 sanity check 的定位，不重复讲二阶和四阶。")

    # 8
    slide = new_slide(prs, 8, "实验二：Neumann/mixed 边界", "ghost-point、DCT-I 与 zero mode")
    add_multiline(slide, [
        "Neumann 给定法向导数，不是直接给函数值",
        "ghost-point 会改变边界行结构，需要 DCT-I 对称化处理",
        "Pure Neumann Poisson 存在常数零模态，需要 weighted mean 归一化",
        "本实验不声明 FFT9 的 Neumann/mixed 四阶支持",
    ], 0.95, 1.55, 11.25, 2.05, size=20)
    add_formula(slide, "误差收敛 + 边界 residual + zero-mode 诊断", 1.6, 4.2, 10.2, 0.80, size=22)
    add_takeaway(slide, "实验二的价值是边界闭合，而不是重复五点二阶本身。")
    note(slide, "解释为什么 Neumann 更麻烦：导数边界、ghost point、零模态。这页为图 6 铺垫。")

    # 9
    slide = new_slide(prs, 9, "实验二：Neumann/mixed summary", "误差、边界约束与零模态诊断")
    img = FIG_DIR / "exp03_neumann_mixed_summary.png"
    if not add_image_fit(slide, img, 0.70, 1.18, 11.95, 5.28):
        missing.append(img.name)
    add_takeaway(slide, "误差收敛、边界残差和 weighted mean 诊断同时闭合。")
    note(slide, "按 panel (a)(b)(c) 顺序讲：收敛、边界 residual、zero-mode。说明 CR/FACR-like 完整结果在 CSV 中，图里只画 FA 避免冗余。")

    # 10
    slide = new_slide(prs, 10, "实验三：精度--成本对比", "论文题目中“对比研究”的核心证据")
    add_multiline(slide, [
        "同一个 Dirichlet manufactured problem 下比较 error 与 solve time",
        "FA/CR/FACR-like：五点二阶直接法",
        "FFT9：Dirichlet 四阶直接法",
        "GMRES30：无预处理重启 GMRES，误差含离散误差与代数残差",
    ], 0.95, 1.45, 11.25, 2.1, size=20)
    add_text(slide, "Timing scope：direct solvers 为 solver call；GMRES 为 core solve，matrix/RHS setup excluded。", 1.0, 4.18, 11.2, 0.42, size=17, bold=True, color=BLUE)
    add_text(slide, "因此该图展示当前实现下的整体基准趋势，不是严格 kernel-to-kernel benchmark。", 1.0, 4.70, 11.2, 0.42, size=17, color=MUTED)
    add_takeaway(slide, "这组实验直接回答：更快和更准能否同时实现。")
    note(slide, "主动说明 timing scope。GMRES setup 被排除，所以比较并没有刻意偏向 FFT 方法。")

    # 11
    slide = new_slide(prs, 11, "实验三：error-time 图", "越靠左下表示越快且越准")
    img = FIG_DIR / "exp06_accuracy_cost_error_time.png"
    if not add_image_fit(slide, img, 0.65, 1.15, 12.05, 5.34):
        missing.append(img.name)
    add_takeaway(slide, "横轴时间、纵轴误差，左下最优；FFT9 优势来自四阶 Dirichlet 离散，GMRES30 是无预处理五点基线。")
    note(slide, "先讲读图规则：横轴是 median solve time，纵轴是 L_infty error，都是 log scale，所以越靠左下越快且越准。左 panel 是 sigma=0 的 Poisson，右 panel 是 sigma=10 的 modified Helmholtz；这不是 true Helmholtz near-resonance 图。FA/CR/FACR-like 和 GMRES30 都对应五点二阶离散，GMRES30 即使代数残差很小，PDE 误差也受五点离散误差限制；FFT9 求解 Dirichlet 九点紧致四阶系统，不是同一个五点系统更准。timing scope 也要主动说：direct solvers 是 solver call，GMRES core solve 排除了 matrix/RHS setup，因此这不是严格 kernel benchmark，且口径上并未刻意偏向 FFT。结论限定为当前规则 Dirichlet 光滑问题和 unpreconditioned GMRES(30) 基线。")

    # 12
    slide = new_slide(prs, 12, "实验三：time scaling 图", "当前实现的规模增长趋势")
    img = FIG_DIR / "exp06_time_scaling.png"
    if not add_image_fit(slide, img, 0.72, 1.20, 11.9, 5.25):
        missing.append(img.name)
    add_takeaway(slide, "FFT 直接法呈 O(N² log N) 趋势；FACR-like 当前实现不声称 O(N² log log N)。")
    note(slide, "讲清楚经典 FACR 与本文 FACR-like 实现的区别。不要把理论最优复杂度写成当前实现结果。")

    # 13
    slide = new_slide(prs, 13, "实验四：Helmholtz 的关键是分母", "modified 稳定，true 可能穿过离散谱")
    add_formula(slide, "dₚ,q = λₚ,q + σ", 0.95, 1.55, 4.5, 0.72, size=25)
    add_formula(slide, "true: dₚ,q = λₚ,q - κ²", 6.85, 1.55, 5.0, 0.72, size=25)
    add_multiline(slide, [
        "modified Helmholtz：σ=+κ²，分母全正",
        "true-away：分母出现正负分裂，但最近模态仍离 0 不近",
        "true-near：目标模态分母接近 0，触发小分母风险",
        "后续 GMRES 行为与这一谱结构相对应",
    ], 1.0, 3.15, 11.1, 1.65, size=20)
    add_takeaway(slide, "实验四把 Helmholtz 类型差异转化为 dₚ,q 的可视化与残差证据。")
    note(slide, "这页连接理论和图 10。modified 与 true 的根本差别是 sigma 的符号使分母性质改变。")

    # 14
    slide = new_slide(prs, 14, "图 10：小分母风险图", "低频 zoom risk map + sorted |d|")
    img = FIG_DIR / "exp07_spectral_denominator_heatmaps.png"
    if not add_image_fit(slide, img, 0.62, 1.10, 12.1, 5.42):
        missing.append(img.name)
    add_takeaway(slide, "上排看热点位置，下排看数量级：modified≈10²，true-away≈10¹，true-near 的 (2,3)/(3,2) 降至 10⁻²。")
    note(slide, "这页是 Helmholtz 实验的谱结构地基。图中画的是 R=-log10|d|，不是解、误差或条件数；R<0 只表示 |d|>1、分母大、风险低。上排低频窗口告诉我们小分母热点在哪里，下排 sorted |d| 直接比较数量级：modified 的最小分母约 1.197e2，true-away 约 1.480e1，说明 sign-changing 但还不是 near-resonance；true-near 的前两个最小分母降到 1.000e-2，并且对应 (2,3)/(3,2)。这两个点并列是因为 lambda_23^h=lambda_32^h，说明危险集中在目标简并模态，而不是整个频域都危险。")

    # 15
    slide = new_slide(prs, 15, "Condition check", "spectral denominator indicator 与 cond₂(A)")
    img = FIG_DIR / "exp04_condition_check.png"
    if not add_image_fit(slide, img, 0.86, 1.25, 11.55, 5.12):
        missing.append(img.name)
    add_takeaway(slide, "Dirichlet 五点系统中，正交 DST-I 对角化使谱分母指标与 cond₂(A) 一致。")
    note(slide, "讲清这是当前 Dirichlet 五点系统的数值校验。不要外推到 Neumann 原始矩阵、mixed 边界或一般 non-normal 系统。")

    # 16
    slide = new_slide(prs, 16, "GMRES residual history", "比最终迭代数更能说明停滞过程")
    img = FIG_DIR / "exp04_gmres_history.png"
    if not add_image_fit(slide, img, 0.70, 1.18, 11.95, 5.30):
        missing.append(img.name)
    add_takeaway(slide, "absolute residual 是停止准则，relative residual 只用于归一化比较。")
    note(slide, "强调 residual history 的价值：它显示残差是否真的达到 tol，还是到上限仍停滞。")

    # 17
    slide = new_slide(prs, 17, "实验五：near-resonance 理论公式", "从小分母到目标模态放大")
    add_formula(slide, "ûₚ,q = f̂ₚ,q / (λₚ,q - κ²)", 1.0, 1.45, 5.8, 0.78, size=23)
    add_formula(slide, "κ² = λ_target + δ", 7.05, 1.45, 4.9, 0.78, size=23)
    add_formula(slide, "|û_target| ≈ |f̂_target| / |δ|", 2.0, 2.65, 9.2, 0.78, size=24)
    add_multiline(slide, [
        "RHS 使用 off-center Gaussian，避免纯 eigenfunction 的特殊收敛行为",
        "target sets 包括 (1,1)、(2,3)/(3,2)、(3,3)",
        "结论限定于当前 Dirichlet 五点离散谱和无预处理 GMRES30",
    ], 1.05, 4.18, 11.2, 1.2, size=19)
    add_takeaway(slide, "实验五把“小分母”变成可量化的模态投影公式。")
    note(slide, "这页是实验五的理论中心。一定讲出目标模态分母等于 -delta，因此幅值按 1/|delta| 放大。")

    # 18
    slide = new_slide(prs, 18, "Baseline near-resonance scan", "同一模态对下的 detuning 对比")
    img = FIG_DIR / "exp05_near_resonance_summary.png"
    if not add_image_fit(slide, img, 0.70, 1.18, 11.95, 5.30):
        missing.append(img.name)
    add_takeaway(slide, "panel (a) 验证 1/|delta| 放大；raw fields 显示 max |u| 约 0.56→564；GMRES 未达到绝对残差容差。")
    note(slide, "这张图按四块讲：panel (a) 是核心，|u_hat_(2,3)| 与 C/|delta| 参考线一致，说明目标模态按小分母公式放大；panel (b) 画的是 final absolute residual，红叉表示 capped/not converged，停止准则是绝对残差容差。panel (c)(d) 是同一模态对 (2,3)/(3,2) 下 delta=1e-1 和 1e-4 的 raw fields，二者使用独立 colorbar，所以颜色不能直接比较；要看标题中的 max |u|，大约从 0.56 放大到 564，约 1000 倍。结论只限当前 Dirichlet 五点离散、当前 RHS 和无预处理 GMRES(30)，不外推到预处理 Krylov 方法。")

    # 19
    slide = new_slide(prs, 19, "Multimode resonance summary", "1/|delta| 放大与 GMRES capped 汇总")
    img = FIG_DIR / "exp05_multimode_resonance_summary.png"
    if not add_image_fit(slide, img, 0.64, 1.12, 12.05, 5.38):
        missing.append(img.name)
    add_takeaway(slide, "near-resonance 不是单一例子：简并对系数按 1/|delta| 放大，三类目标的能量/形状指标接近 1，GMRES30 全部 capped。")
    note(slide, "这页用于说明 near-resonance 机制不是只对 (2,3)/(3,2) 一个例子成立。当前图的 panel (a) 专门展开简并对内部两个系数 |u_hat_(2,3)| 和 |u_hat_(3,2)|；二者高度可以不同，因为 Gaussian RHS 在两个模态上的投影不同，但斜率都符合 1/|delta|。panel (b) 是 target subspace energy fraction，panel (c) 是 full solution 与 dominant projection 的 shape correlation；它们接近 1 不是没有信息，而是说明在当前 detuning 区间内目标子空间支配性已经建立。panel (d) 只汇总 capped=1001 的终点状态，停滞过程要接下一张 residual history；结论仍限定为当前 Dirichlet 五点离散、当前 RHS 和 unpreconditioned GMRES(30)。")

    # 20
    slide = new_slide(prs, 20, "Dominant projection", "从频域结论回到物理空间解场")
    img = FIG_DIR / "exp05_dominant_mode_projection.png"
    if not add_image_fit(slide, img, 0.58, 1.24, 9.05, 5.20):
        missing.append(img.name)
    add_text(slide, "量级读法", 9.95, 1.30, 2.65, 0.28, size=17, bold=True, color=TITLE)
    add_multiline(slide, [
        "full max ≈ 6.81×10¹",
        "projection max ≈ 6.81×10¹",
        "difference max ≈ 6.89×10⁻⁴",
        "相对量级 ≈ 1×10⁻⁵",
    ], 9.95, 1.72, 2.75, 1.75, size=14.5)
    add_text(slide, "右图使用独立色标：颜色明显不代表差值大，要看 max |u| 数量级。", 9.95, 3.72, 2.70, 0.95, size=14.5, color=WARN)
    add_text(slide, "差值来自非目标模态的有限贡献，不是 projection 失败。", 9.95, 4.92, 2.70, 0.72, size=14.5, color=MUTED)
    add_takeaway(slide, "完整解与目标模态投影只差约 10⁻⁵ 相对量级，near-resonance 解场由 (2,3)/(3,2) 主导。")
    note(slide, "这一页要讲清楚：左图是完整解，中图是只保留 (2,3)/(3,2) 后反变换得到的 dominant projection，右图是差值。full 和 projection 最大幅值都是约 6.81e1，difference 只有约 6.89e-4，差了五个数量级。右图颜色明显是因为用了独立色标，不代表误差大；差值不是机器零，因为 Gaussian RHS 对非目标模态也有小投影。")

    # 21
    slide = new_slide(prs, 21, "Resonance GMRES history", "capped 结果背后的残差过程")
    img = FIG_DIR / "exp05_resonance_gmres_history.png"
    if not add_image_fit(slide, img, 0.70, 1.18, 11.95, 5.30):
        missing.append(img.name)
    add_takeaway(slide, "not converged 不是单一标签：曲线先有限下降后平台，末端约 1.8e-1，仍远高于 1e-10；capped=1001 表示达到迭代上限。")
    note(slide, "这一页按四步讲。第一，左图是 absolute residual，它对应真正的停止准则 1e-10；右图是 relative residual，只用于归一化比较。第二，三条曲线并不是一开始完全失败，而是先有限下降，然后进入平台，末端约 1.8e-1，远高于容差。第三，红叉和 capped=1001 表示 max_iter=1000 后仍未满足绝对残差容差，不是第 1001 步成功收敛。第四，结论限定在当前 Dirichlet 五点离散、目标模态 (2,3)/(3,2)、当前 RHS 和 unpreconditioned GMRES(30)，不能外推到预处理 GMRES、MINRES 或 shifted Laplacian。")

    # 22
    slide = new_slide(prs, 22, "五个实验对应五个结论", "第六章的证据链收束")
    add_table(slide, [
        ["实验", "回答的问题", "核心结论"],
        ["Dirichlet 收敛", "格式阶数是否达标", "五点二阶，FFT9 四阶"],
        ["Neumann/mixed", "边界与零模态是否闭合", "residual 与 weighted mean 诊断闭合"],
        ["精度--成本", "FFT 与 GMRES 如何比较", "FFT9 在光滑 Dirichlet 下精度--时间更优"],
        ["谱结构+GMRES", "true Helmholtz 为什么难", "小分母与残差停滞对应"],
        ["near-resonance", "具体哪个模态被放大", "目标模态遵循 1/|delta| 放大"],
    ], 0.62, 1.36, 12.1, 4.72, [2.15, 4.25, 5.7], font_size=11.5)
    add_takeaway(slide, "五个实验不是孤立结果，而是从格式验证走向机制解释。")
    note(slide, "总结页不要展开太久。每行一句，突出五个实验的分工。")

    # 23
    slide = new_slide(prs, 23, "结论边界与限制", "主动说明，避免过度声称")
    add_multiline(slide, [
        "GMRES 是无预处理 GMRES30，不代表所有预处理 Krylov 方法",
        "FFT9 仅在 Dirichlet 边界实现和验证",
        "FACR-like 当前实现仍按 O(N² log N) 处理",
        "timing 图不是严格 kernel-to-kernel benchmark",
        "near-resonance 是离散谱实验；condition check 不外推到 non-normal 系统；exp04/exp05 capped 计数约定不同",
    ], 0.95, 1.50, 11.25, 2.65, size=20)
    add_takeaway(slide, "限制越清楚，实验结论越可信。")
    note(slide, "答辩时主动讲限制是加分项。不要说 GMRES 不适合 Helmholtz，也不要说 FFT9 支持 Neumann/mixed。")

    # 24
    slide = new_slide(prs, 24, "答辩口径：常见问题与短答", "限定条件 + 证据来源 + 不外推")
    add_table(slide, [
        ["问题", "短答"],
        ["exp00 为什么不算核心实验？", "只验证同一离散系统的求解路径一致。"],
        ["FFT9 优势来自哪里？", "Dirichlet 四阶离散误差，而非迭代残差。"],
        ["GMRES 结论能外推吗？", "不能；这里只是无预处理 GMRES30 基准。"],
        ["risk map 是条件数图吗？", "不是；它只可视化 small-denominator risk。"],
        ["图 16(b)(c) 为什么近似平线？", "值已接近 1，说明目标模态支配了解。"],
        ["GMRES 计时是否偏向 FFT？", "GMRES setup 被排除，计时更宽容；FFT 仍占优。"],
    ], 0.62, 1.28, 12.1, 5.05, [4.65, 7.45], font_size=10.7)
    add_takeaway(slide, "答辩时围绕“问题--证据--限制”回答。")
    note(slide, "最后一页作为 Q&A 备忘。回答要短，不要自发扩展到预处理方法、连续谱极限或未实现边界。")

    prs.save(PPTX_PATH)
    return missing, note_count


def validate(note_count: int) -> dict[str, object]:
    result: dict[str, object] = {
        "exists": PPTX_PATH.exists(),
        "size_bytes": PPTX_PATH.stat().st_size if PPTX_PATH.exists() else 0,
        "slide_xml_count": 0,
        "media_count": 0,
        "content_types": False,
        "remarks_count": 0,
        "python_pptx_open": False,
        "python_pptx_slides": 0,
        "expected_remarks": note_count,
    }
    if not PPTX_PATH.exists():
        return result
    with zipfile.ZipFile(PPTX_PATH, "r") as zf:
        names = zf.namelist()
        result["slide_xml_count"] = len([n for n in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)])
        result["media_count"] = len([n for n in names if n.startswith("ppt/media/")])
        result["content_types"] = "[Content_Types].xml" in names
        result["remarks_count"] = sum(
            zf.read(n).decode("utf-8", errors="ignore").count("SPEAKER_NOTES:")
            for n in names
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)
        )
    try:
        prs = Presentation(str(PPTX_PATH))
        result["python_pptx_open"] = True
        result["python_pptx_slides"] = len(prs.slides)
    except Exception as exc:
        result["python_pptx_error"] = repr(exc)
    return result


def write_report(missing_figs: list[str], validation: dict[str, object]) -> None:
    missing_csv = [c.filename for c in CSVS if not (CSV_DIR / c.filename).exists()]
    status = (
        "READY TO USE"
        if not missing_figs and not missing_csv and validation.get("slide_xml_count") == 24 and validation.get("remarks_count") == 24
        else "READY WITH MISSING ASSETS"
    )
    lines = [
        "# CH6 PPT Build Report",
        "",
        "## 1. Outputs",
        f"- PPTX: `{rel(PPTX_PATH)}`",
        f"- Build script: `{rel(Path(__file__))}`",
        f"- HTML companion: `{rel(HTML_PATH)}`",
        "",
        "## 2. Build Scope",
        "- Rebuilt a 16:9 editable PPTX as a graduate-thesis experiment report deck.",
        f"- Dependency: `python-pptx {pptx.__version__}`.",
        "- Used existing figures from `thesis/figures/` only.",
        "- Did not modify thesis text, experiment scripts, CSV data, or figure generation code.",
        "- Did not add experiments or redraw figures.",
        "- Speaker guidance is embedded as off-slide remark text boxes (`SPEAKER_NOTES:`).",
        "",
        "## 3. Figures Used",
    ]
    for fig in FIGURES:
        path = FIG_DIR / fig.filename
        marker = "OK" if path.exists() else "MISSING"
        lines.append(f"- Slide {fig.slide}: `{rel(path)}` - {fig.purpose} [{marker}]")
    lines += ["", "## 4. CSV Files Checked"]
    for c in CSVS:
        path = CSV_DIR / c.filename
        marker = "OK" if path.exists() else "MISSING"
        lines.append(f"- `{rel(path)}` - {c.purpose} [{marker}]")
    lines += [
        "",
        "## 5. Missing Assets",
        f"- Missing figures: {len(missing_figs)}",
    ]
    for name in missing_figs:
        lines.append(f"  - `{name}`")
    lines.append(f"- Missing CSV files: {len(missing_csv)}")
    for name in missing_csv:
        lines.append(f"  - `{name}`")
    lines += [
        "",
        "## 6. Structural Validation",
        f"- PPTX exists: {validation.get('exists')}",
        f"- PPTX size: {validation.get('size_bytes')} bytes",
        f"- Slide XML count: {validation.get('slide_xml_count')}",
        f"- python-pptx slide count: {validation.get('python_pptx_slides')}",
        f"- Media resources: {validation.get('media_count')}",
        f"- `[Content_Types].xml` present: {validation.get('content_types')}",
        f"- Off-slide remark count: {validation.get('remarks_count')}",
        f"- Expected remark count: {validation.get('expected_remarks')}",
        f"- Openable by python-pptx: {validation.get('python_pptx_open')}",
        "- Visual rendering/export preview: not run; structural PPTX validation completed.",
        "",
        "## 7. Slide Count",
        "- Total slides: 24",
        "- Format: 16:9 widescreen",
        "- Audience: graduate thesis defense / experiment report",
        "",
        "## 8. Final Status",
        status,
        "",
    ]
    BUILD_REPORT.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    missing_figs, note_count = build_deck()
    validation = validate(note_count)
    write_report(missing_figs, validation)
    print(f"Wrote {PPTX_PATH}")
    print(f"Wrote {BUILD_REPORT}")
    print(f"slides={validation.get('slide_xml_count')} remarks={validation.get('remarks_count')} missing_figures={len(missing_figs)}")


if __name__ == "__main__":
    main()
